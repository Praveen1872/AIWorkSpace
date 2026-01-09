import streamlit as st
import firebase_admin
from firebase_admin import credentials, firestore
from pptx import Presentation
import io, json, re
import os 
from google import genai
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor



st.set_page_config(page_title="Slide Architect Pro", layout="wide")
# st.markdown("""
# <style>
#     .stApp { 
#         background-color: #FAF8F7; 
#         color: #1A1A1A; 
#         font-family: 'Inter', sans-serif;
#     }

#     .slide-stage {
#         width: 700px;
#         height: 500px;
#         background-color: #FFFFFF;
#         margin: 20px auto;
#         padding: 50px;
#         border-radius: 8px;
#         border: 1px solid #E0DEDD;
#         box-shadow: 0px 10px 30px rgba(0, 0, 0, 0.08);
#         display: flex;
#         flex-direction: column;
#         overflow: hidden;
#         position: relative;
#     }
# .slide-stage {
#         background-color: #ffffff;
#         border: 1px solid #e2e8f0;
#         border-radius: 12px;
#         padding: 40px;
       
#         box-shadow: 0 10px 15px -3px rgba(0,0,0,0.1);
#         margin-bottom: 20px;
#     }
#     .slide-title {
#         font-size: 42px;
#         font-weight: 800;
#         color: #FF4520;
#         margin-bottom: 25px;
#         line-height: 1.2;
#     }

#     .slide-content {
#         font-size: 26px;
#         line-height: 1.6;
#         color: #2D2D2D;
#         flex-grow: 1;
#     }

#     .slide-point {
#         margin-bottom: 20px;
#         padding-left: 10px;
#     }

#     div.stButton > button { 
#         border-radius: 12px; 
#         background-color: white; 
#         color: black;
#         border: 1px solid #E0DEDD;
#         height: 3.2em;
#         width: 100%;
#         font-weight: 600;
#         transition: all 0.3s cubic-bezier(0.4, 0, 0.2, 1);
#     }

#     div.stButton > button:hover {
#         background-color: #FF4520;
#         color: white !important;
#         transform: translateY(-2px);
#         box-shadow: 0px 6px 15px rgba(255, 69, 32, 0.3);
#         border-color: #FF4520;
#     }

#     [data-testid="column"] {
#         display: flex;
#         align-items: center;
#         justify-content: center;
#         padding: 0px 5px !important;
#     }
# </style>
# """, unsafe_allow_html=True)



is_logged_in = st.session_state.get('logged_in', False)
if not is_logged_in:
    st.switch_page("pages/login.py")

if "ppt_data" not in st.session_state:
    st.session_state.ppt_data = []

if "slide_styles" not in st.session_state:
    st.session_state.slide_styles = {}

if "current_slide_idx" not in st.session_state:
    st.session_state.current_slide_idx = 0

if st.session_state.current_slide_idx >= len(st.session_state.ppt_data):
    st.session_state.current_slide_idx = 0


user_uid = st.session_state.get("user_uid")



def initialize_firebase():
    if not firebase_admin._apps:
        try:
            cred = credentials.Certificate({
                "type": "service_account",
                "project_id": os.getenv("FIREBASE_PROJECT_ID"),
                "private_key": os.getenv("FIREBASE_PRIVATE_KEY").replace("\\n", "\n"),
                "client_email": os.getenv("FIREBASE_CLIENT_EMAIL"),

                # 🔑 REQUIRED EXTRA FIELDS
                "token_uri": "https://oauth2.googleapis.com/token",
                "auth_uri": "https://accounts.google.com/o/oauth2/auth",
            })

            firebase_admin.initialize_app(cred)

        except Exception as e:
            raise RuntimeError(f"Firebase Initialization Failed: {e}")

    return firestore.client()


firestore_db = initialize_firebase()
user_uid = st.session_state.get("user_uid")

ppt_col = (
    firestore_db
    .collection("users")
    .document(user_uid)
    .collection("documents")
    .document("ppt")
    .collection("items")
)



h_cols = st.columns([2, 0.9, 0.9, 0.9, 1.5, 0.8, 1], vertical_alignment="center")
with h_cols[0]: 
    st.markdown("<h3 style='margin:0;'>🚀 AI Mentor</h3>", unsafe_allow_html=True)

with h_cols[1]: 
    if st.button("PPT 🖼️", use_container_width=True,type="primary"): st.switch_page("pages/ppt_editor.py")
with h_cols[2]: 
    if st.button("Word 📝", use_container_width=True): st.switch_page("pages/word_editor.py")
with h_cols[3]:
    
    if st.button("Notes 📓", use_container_width=True): st.switch_page("pages/note.py")
with h_cols[4]: 
    
    if st.button("Summarizer 📝", use_container_width=True): st.switch_page("pages/Summarizer.py")

with h_cols[6]:
    if st.button("Logout 🚪", use_container_width=True):
        st.session_state.logged_in = False
        st.switch_page("AIMentor.py")


st.markdown("<hr style='margin:0 0 20px 0;'>", unsafe_allow_html=True)

# st.markdown("""
#     <style>
    
#     .slide-title { color: #1e293b; font-size: 32px; font-weight: 800; margin-bottom: 25px; border-bottom: 3px solid #3b82f6; width: 100%; padding-bottom: 10px; }
#     .col-container { display: flex; gap: 30px; }
#     .content-col { flex: 1; }
#     .slide-point { font-size: 18px; margin-bottom: 12px; color: #475569; line-height: 1.6; }
#     .mentor-box { background-color: #f0fdf4; border-left: 5px solid #22c55e; padding: 15px; border-radius: 5px; margin-top: 10px; }
#     </style>
# """, unsafe_allow_html=True)

def clean_text(text):
    return re.sub(r'\*\*|#+', '', str(text)).strip()



import os

API_KEY = os.getenv("GEMINI_API_KEY")

client = genai.Client(api_key=API_KEY)
def call_ai_architect(prompt, current_data=None, active_idx=None):
    system_instr = (
        "You are an Academic Slide Architect.\n"
        "Create slides with clear titles and concise bullet points.\n"
        "Each slide MUST have a MAXIMUM of 4 bullet points.\n\n"
        "STRICT RULES:\n"
        "- Return ONLY valid JSON\n"
        "- Do NOT include HTML, markdown, bullets, symbols, or numbering\n"
        "- Points must be plain text sentences only\n\n"
        "Format:\n"
        "{'slides': [{'title': 'Slide title', 'points': ['Point one', 'Point two']}], "
        "'mentor_advice': '...'}"
    )

    if current_data:
        focus = f"Slide {active_idx + 1}" if active_idx is not None else "the whole deck"
        system_instr = (
            f"Context: {json.dumps(current_data)}\n"
            f"Update {focus} based on: {prompt}\n\n"
            + system_instr
        )

    model_list = ["gemini-2.5-flash-lite"]

    for model_name in model_list:
        try:
            response = client.models.generate_content(
                model=model_name,
                contents=prompt,
                config={"system_instruction": system_instr},
            )

            # 🔒 Extract JSON only
            match = re.search(r"\{[\s\S]*\}", response.text)
            if not match:
                continue

            res_json = json.loads(match.group(0))
            slides = res_json.get("slides", [])

            clean_slides = []

            for s in slides:
                title = clean_text(s.get("title", ""))

                points = []
                for p in s.get("points", []):
                    p = str(p)
                    p = re.sub(r"<[^>]+>", "", p)           # remove HTML
                    p = re.sub(r"^[•\-–\d\.\s]+", "", p)    # remove bullets/numbers
                    p = clean_text(p)

                    if p:
                        points.append(p)

                if title and points:
                    clean_slides.append({
                        "title": title,
                        "points": points[:4]   # HARD LIMIT
                    })

            if clean_slides:
                return clean_slides, res_json.get("mentor_advice", ""), model_name

        except Exception:
            continue

    return None, None, None


def ai_style_title(instruction):
    """
    Converts natural language into title style JSON
    """
    style_prompt = f"""
You are a slide design assistant.

Convert the instruction below into JSON for a slide title style.

Instruction:
{instruction}

Rules:
- Respond ONLY with valid JSON
- Keys allowed: font, size, weight, color
- Font must be a common PPT font
- Size between 24 and 64
- Weight: 400 (normal) or 800 (bold)
- Color must be hex

Example output:
{{"font":"Montserrat","size":48,"weight":800,"color":"#2563eb"}}
"""

    try:
        response = client.models.generate_content(
            model="gemini-2.5-flash-lite",
            contents=[style_prompt]
        )

        match = re.search(r"\{.*\}", response.text, re.DOTALL)
        if match:
            return json.loads(match.group(0))
    except:
        pass

    return None

def store_ppt_chunks(ppt_doc_ref, slides):
    chunks_col = ppt_doc_ref.collection("chunks")

    for s in slides: 
        chunks_col.add({
            "title": s.get("title", "").strip(),
            "points": s.get("points", []),
            "embedding": []  # placeholder for future RAG
        })


def load_user_ppts():
    docs = ppt_col.order_by(
        "created_at",
        direction=firestore.Query.DESCENDING
    ).stream()

    return [{"id": d.id, **d.to_dict()} for d in docs]

with st.sidebar:
    st.subheader("📂 PPT History")

    ppt_docs = load_user_ppts()
    if ppt_docs:
        title_map = {p["title"]: p["id"] for p in ppt_docs}
        selected = st.selectbox("Past PPTs", title_map.keys())

        if st.button("📂 Load PPT"):
            st.session_state.active_ppt_id = title_map[selected]
            st.rerun()

    if st.button("🗑️ Clear All PPTs"):
        for doc in ppt_col.stream():
            for c in doc.reference.collection("chunks").stream():
                c.reference.delete()
            doc.reference.delete()
        st.session_state.pop("ppt_data", None)
        st.rerun()


if "active_ppt_id" in st.session_state:
    ppt_doc = ppt_col.document(st.session_state.active_ppt_id)
    chunks = ppt_doc.collection("chunks").stream()

    slides = []
    for c in chunks:
        data = c.to_dict()
        slides.append({
            "title": data.get("title", ""),
            "points": data.get("points", []),
            "style": data.get("style", {
                "font": "Arial",
                "size": 42,
                "weight": 800,
                "color": "#1e293b"
            })
        })

    if slides:
        st.session_state.ppt_data = slides



if "slide_styles" not in st.session_state:
    for i in range(len(st.session_state.ppt_data)):
        st.session_state.slide_styles.setdefault(
        i,
        {"font": "Arial", "size": 42, "weight": 800, "color": "#1e293b"}
    )


# Ensure current slide has a style
if "current_slide_idx" in st.session_state:
    idx = st.session_state.current_slide_idx
    if idx not in st.session_state.slide_styles:
        st.session_state.slide_styles[idx] = {
            "font": "Arial",
            "size": 42,
            "weight": 800,
            "color": "#1e293b"
        }


col_stage, col_chat = st.columns([1.8, 1], gap="large")

with col_stage:
    if not st.session_state.get("ppt_data"):
        st.info("👋 Ask the Assistant to generate slides")
    else:
        # Current Slide Data
        data = st.session_state.ppt_data
        idx = st.session_state.current_slide_idx
        active_slide = data[idx]
        
        # Get style (Manual user changes happen here)
        style = st.session_state.slide_styles.get(idx, {"font": "Arial", "size": 42, "weight": 800, "color": "#1e293b"})

        # --- Manual Style UI ---
        with st.expander("🎨 Manual Style Adjustments"):
            c1, c2, c3 = st.columns(3)
            new_font = c1.selectbox("Font", ["Arial", "Poppins", "Montserrat"], index=0)
            new_size = c2.slider("Size", 20, 80, style["size"])
            new_color = c3.color_picker("Color", style["color"])
            
            # Save manual changes back to state
            st.session_state.slide_styles[idx] = {
                "font": new_font, "size": new_size, "weight": 800, "color": new_color
            }

        # --- Render Slide ---
        st.markdown(f"""
            <div style="background: white; padding: 40px; border-radius: 15px; border: 1px solid #ddd;">
                <h1 style="font-family: {style['font']}; font-size: {style['size']}px; color: {style['color']};">
                    {active_slide['title']}
                </h1>
                <ul>
                    {"".join([f"<li>{p}</li>" for p in active_slide['points']])}
                </ul>
            </div>
        """, unsafe_allow_html=True)
ppt_data = st.session_state.get("ppt_data", [])

if len(ppt_data) > 0:
    st.write("### 🎞️ Slide Navigator")

    nav_cols = st.columns(min(len(ppt_data), 10))

    for i in range(len(ppt_data)):
        with nav_cols[i % len(nav_cols)]:
            if st.button(
                str(i + 1),
                key=f"nav_{i}",
                type="primary" if i == st.session_state.current_slide_idx else "secondary"
            ):
                st.session_state.current_slide_idx = i
                st.rerun()

    # ---------- NAVIGATOR ----------
    st.write("### 🎞️ Slide Navigator")
    ppt_data = st.session_state.get("ppt_data", [])
    nav_cols = st.columns(min(len(ppt_data), 10))


    for i in range(len(ppt_data)):
        with nav_cols[i % 10]:
            if st.button(
                str(i + 1),
                key=f"nav_{i}",
                type="primary" if i == st.session_state.current_slide_idx else "secondary"
            ):
                st.session_state.current_slide_idx = i
                st.rerun()

    st.divider()

# ✅ ALWAYS create Presentation ONCE
prs = Presentation()
prs.slide_width = 9144000
prs.slide_height = 6858000

c1, c2 = st.columns(2)

# ---------------- DELETE SLIDE ----------------
with c1:
    if st.button("🗑️ Delete Slide", use_container_width=True):
        if len(st.session_state.ppt_data) > 0:
            st.session_state.ppt_data.pop(st.session_state.current_slide_idx)
        # Shift index back so it doesn't point to a ghost slide
            st.session_state.current_slide_idx = max(0, st.session_state.current_slide_idx - 1)
            st.rerun()

# ---------------- EXPORT PPT ----------------
# ---------------- EXPORT PPT ----------------
with c2:
    if st.button("📥 Download PPTX", use_container_width=True):

        prs = Presentation()
        prs.slide_width = 9144000
        prs.slide_height = 6858000

        for s in st.session_state.ppt_data:
            slide = prs.slides.add_slide(prs.slide_layouts[1])

            # TITLE
            title_shape = slide.shapes.title
            title_shape.text = clean_text(s.get("title", ""))

            style = s.get("style", {})

            for p in title_shape.text_frame.paragraphs:
                for run in p.runs:
                    run.font.name = style.get("font", "Arial")
                    run.font.size = Pt(style.get("size", 42))
                    run.font.bold = style.get("weight", 800) >= 700

                    hex_color = style.get("color", "#000000").lstrip("#")
                    if len(hex_color) == 6:
                        run.font.color.rgb = RGBColor(
                            int(hex_color[0:2], 16),
                            int(hex_color[2:4], 16),
                            int(hex_color[4:6], 16),
                        )

            # CONTENT
            points = s.get("points", [])[:7]
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = "\n".join(clean_text(p) for p in points)

        buf = io.BytesIO()
        prs.save(buf)

        st.download_button(
            "⬇ Download PPTX",
            data=buf.getvalue(),
            file_name="presentation.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            use_container_width=True,
        )


with col_chat:
    st.title("AI Assistant")
    edit_mode = st.toggle("🎯 Edit ONLY Active Slide", value=False)
    
    if "chat_history" not in st.session_state:
        st.session_state.chat_history = []

    chat_box = st.container(height=450)
    for m in st.session_state.chat_history:
        with chat_box.chat_message(m["role"]):
            st.write(m["content"])
    
            if "advice" in m:
                st.markdown(f'<div class="mentor-box">💡 {m["advice"]}</div>', unsafe_allow_html=True)
    st.markdown("### 🎨 AI Title Styling (Active Slide)")
    ppt_prompt = st.chat_input("Enter PPT topic")
ai_style_cmd = st.text_input(
    "Describe title style (e.g., modern bold blue)",
    key="ai_title_style"
)

if st.button("✨ Apply AI Style"):
    style_update = ai_style_title(ai_style_cmd)

    if style_update:
        st.session_state.slide_styles[
            st.session_state.current_slide_idx
        ].update(style_update)

        st.success("Title style updated")
        st.rerun()
    else:
        st.warning("Could not understand style command")

if user_in := ppt_prompt:
    st.session_state.chat_history.append({"role": "user", "content": user_in})

    with st.spinner("Architecting ..."):
        idx = st.session_state.current_slide_idx if edit_mode else None
        new_slides, advice, model_name = call_ai_architect(user_in, st.session_state.ppt_data, idx)

        if new_slides:
            if not edit_mode:
                st.session_state.ppt_data = new_slides
                # Reset styles for a brand new presentation
                st.session_state.slide_styles = {} 
            else:
                st.session_state.ppt_data[st.session_state.current_slide_idx] = new_slides[0]

            # ✅ Ensure every slide has a style entry
            for i in range(len(st.session_state.ppt_data)):
                if i not in st.session_state.slide_styles:
                    st.session_state.slide_styles[i] = {
                        "font": "Arial", "size": 42, "weight": 800, "color": "#1e293b"
                    }

            st.session_state.chat_history.append({"role": "assistant", "content": "Updated!", "advice": advice})
            st.rerun()
        for i in range(len(st.session_state.ppt_data)):
            st.session_state.slide_styles.setdefault(
        i,
        {"font": "Arial", "size": 42, "weight": 800, "color": "#1e293b"}
    )
