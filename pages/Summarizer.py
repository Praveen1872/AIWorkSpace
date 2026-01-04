import streamlit as st
import io, re
from google import genai
import PyPDF2
from docx import Document
from pptx import Presentation
from reportlab.platypus import SimpleDocTemplate, Paragraph
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib.pagesizes import A4
import firebase_admin
from firebase_admin import credentials, firestore

st.set_page_config(
    page_title="Summarizer Lab",
    layout="wide",
    initial_sidebar_state="collapsed"
)


st.markdown("""
<style>
.stApp {
    background-color: #FAF8F7;
    color: #1A1A1A;
    font-family: 'Inter', sans-serif;
}
div.stButton > button {
    border-radius: 12px;
    background-color: white;
    color: black;
    border: none;
    height: 3.2em;
    width: 80%;
    font-weight: 600;
    transition: all 0.3s ease;
}
div.stButton > button:hover {
    background-color: #FF4520;
    color: white !important;
}
hr {
    border-top: 1px solid #E0DEDD;
}
    .summary-scroll {
    max-height: 420px;
    overflow-y: auto;
    padding: 15px;
    background-color: #FFFFFF;
    border-radius: 10px;
    border: 1px solid #E0DEDD;
}
</style>
""", unsafe_allow_html=True)


is_logged_in = st.session_state.get('logged_in', False)
if not is_logged_in:
    st.switch_page("pages/login.py")


user_uid = st.session_state.get("user_uid")

def initialize_firebase():
    if not firebase_admin._apps:
        creds = dict(st.secrets["firebase_credentials"])
        creds["private_key"] = creds["private_key"].replace("\\n", "\n")
        cred = credentials.Certificate(creds)
        firebase_admin.initialize_app(cred)

initialize_firebase()
user_uid = st.session_state.get("user_uid")
firestore_db = firestore.client()
summary_col = (
    firestore_db
    .collection("users")
    .document(user_uid)
    .collection("documents")
    .document("summarized")
    .collection("items")
)

def shorten_title(title, max_len=40):
    if len(title) <= max_len:
        return title
    return title[:max_len].rsplit(" ", 1)[0] + "..."

h_cols = st.columns([2, 0.9, 0.9, 0.9, 1.5, 0.8, 1], vertical_alignment="center")
with h_cols[0]: 
    st.markdown("<h3 style='margin:0;'>🚀 AI Mentor</h3>", unsafe_allow_html=True)

with h_cols[1]: 
    if st.button("PPT 🖼️", use_container_width=True): st.switch_page("pages/ppt_editor.py")
with h_cols[2]: 
    if st.button("Word 📝", use_container_width=True): st.switch_page("pages/word_editor.py")
with h_cols[3]: 
    
    if st.button("Notes 📓", use_container_width=True): st.switch_page("pages/note.py")
with h_cols[4]: 
    
    if st.button("Summarizer📝", use_container_width=True): st.switch_page("pages/Summarizer.py")

with h_cols[6]:
    if st.button("Logout 🚪", use_container_width=True):
        st.session_state.logged_in = False
        st.switch_page("AIMentor.py")

st.markdown("<hr style='margin:0 0 20px 0; border-top: 1px solid #E0DEDD;'>", unsafe_allow_html=True)
# -------------------- Session State Init --------------------
if "active_context" not in st.session_state:
    st.session_state.active_context = ""

if "active_filename" not in st.session_state:
    st.session_state.active_filename = ""

if "summary_output" not in st.session_state:
    st.session_state.summary_output = ""

if "assistant_chat" not in st.session_state:
    st.session_state.assistant_chat = []

# -------------------- AI Client --------------------
client = genai.Client(api_key=st.secrets["GEMINI_API_KEY"])

# -------------------- Helpers --------------------
def extract_text_from_any(uploaded_file):
    name = uploaded_file.name.lower()
    text = ""

    if name.endswith(".pdf"):
        reader = PyPDF2.PdfReader(uploaded_file)
        text = "\n".join([p.extract_text() or "" for p in reader.pages])

    elif name.endswith(".docx"):
        doc = Document(uploaded_file)
        text = "\n".join([p.text for p in doc.paragraphs])

    elif name.endswith(".pptx"):
        prs = Presentation(uploaded_file)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + " "

    return text.strip()

def call_research_ai(prompt, context):
    system = (
        "You are an academic research assistant. "
        "Answer clearly and concisely based ONLY on the provided context."
    )

    response = client.models.generate_content(
        model="gemini-2.5-flash-tts",
        contents=prompt,
        config={
            "system_instruction": f"{system}\n\nContext:\n{context[:18000]}"
        }
    )
    return response.text

def generate_summary_pdf(summary_text):
    buffer = io.BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=A4)
    styles = getSampleStyleSheet()
    story = []

    for line in summary_text.split("\n"):
        story.append(Paragraph(line, styles["Normal"]))

    doc.build(story)
    buffer.seek(0)
    return buffer

col_left, col_right = st.columns([2, 1], gap="large")


with col_left:
    st.title("📑 Summarizer Lab")

    # Attach file
    with st.container(border=True):
        file = st.file_uploader(
            "📎 Attach PDF / DOCX / PPTX",
            type=["pdf", "docx", "pptx"]
        )

        if file and st.button("📥 Attach", use_container_width=True):
            with st.spinner("Reading document..."):
                content = extract_text_from_any(file)
                if content:
                    st.session_state.active_context = content
                    st.session_state.active_filename = file.name
                    st.session_state.summary_output = ""
                    st.success(f"Attached: {file.name}")

    # Generate summary
    if st.session_state.active_context:
        if st.button("✨ Generate Summary", type="primary", use_container_width=True):
            with st.spinner("Generating summary..."):
                st.session_state.summary_output = call_research_ai(
                    "Provide a structured summary with Abstract, Key Findings, and Technical Implications.",
                    st.session_state.active_context
                )
if st.session_state.summary_output:
    st.markdown("### 📄 Summary")

    st.markdown(
        f"""
        <div class="summary-scroll">
            {st.session_state.summary_output}
        </div>
        """,
        unsafe_allow_html=True
    )

    pdf_buffer = generate_summary_pdf(st.session_state.summary_output)

    c1, c2 = st.columns(2)

    with c1:
        st.download_button(
            "📥 Download PDF",
            data=pdf_buffer,
            file_name="summary.pdf",
            mime="application/pdf",
            use_container_width=True
        )

    with c2:
        if st.button("🧹 Clear Summary", use_container_width=True):
            st.session_state.summary_output = ""
            st.session_state.active_context = ""
            st.session_state.active_filename = ""
            st.session_state.assistant_chat = []
            st.rerun()


# ==================== RIGHT: AI ASSISTANT ====================
with col_right:
    st.title("🤖 AI Assistant")
    st.caption("Instant doubt clarification from the attached document")

    if not st.session_state.active_context:
        st.info("📎 Attach a document from the left panel to begin.")
    else:
        chat_box = st.container(height=420)

        with chat_box:
            for msg in st.session_state.assistant_chat:
                with st.chat_message(msg["role"]):
                    st.markdown(msg["content"])

        question = st.chat_input("Ask a question about this document")

        if question:
            st.session_state.assistant_chat.append({
                "role": "user",
                "content": question
            })

            with st.spinner("Thinking..."):
                answer = call_research_ai(
                    question,
                    st.session_state.active_context
                )

            st.session_state.assistant_chat.append({
                "role": "assistant",
                "content": answer
            })

            st.rerun()

def store_summary_chunks(summary_doc_ref, summary_text):
    chunks_col = summary_doc_ref.collection("chunks")

    # Split by paragraphs and clean
    paragraphs = [
        p.strip()
        for p in summary_text.split("\n")
        if len(p.strip()) >= 50   # avoid weak/noisy chunks
    ]

    for idx, p in enumerate(paragraphs):
        chunks_col.add({
            "text": p,
            "embedding": [],      # placeholder for future RAG
            "order": idx          # preserves original order (IMPORTANT later)
        })
def load_user_summaries():
    docs = (
        summary_col
        .order_by("created_at", direction=firestore.Query.DESCENDING)
        .stream()
    )

    summaries = []
    for d in docs:
        data = d.to_dict() or {}
        summaries.append({
            "id": d.id,
            "title": data.get("title", "Untitled Summary"),
            "created_at": data.get("created_at"),
            "source": data.get("source", "summarizer")
        })

    return summaries
with st.sidebar:
    st.subheader("📂 Summary History")

    summaries = load_user_summaries()

    if summaries:
        title_map = {
            f"{i+1}. {shorten_title(s['title'])}": s["id"]
            for i, s in enumerate(summaries)
        }

        selected = st.selectbox(
            "Past Summaries",
            title_map.keys()
        )

        if st.button("📂 Load Summary"):
            st.session_state.active_summary_id = title_map[selected]
            st.rerun()

    if st.button("🗑️ Clear All Summaries"):
        for doc in summary_col.stream():
            for c in doc.reference.collection("chunks").stream():
                c.reference.delete()
            doc.reference.delete()

        st.session_state.pop("summary_text", None)
        st.rerun()

if "active_summary_id" in st.session_state:
    doc_ref = summary_col.document(st.session_state.active_summary_id)
    chunks = doc_ref.collection("chunks").stream()

    summary_text = "\n\n".join(
        c.to_dict()["text"] for c in chunks
    )

    st.session_state.summary_text = summary_text
def generate_title_from_content(content):
    prompt = (
        "Generate a short, clear academic title (max 8 words) "
        "that best represents the following document:\n\n"
        + content[:3000]
    )

    response = client.models.generate_content(
        model="gemini-2.5-flash",
        contents=prompt
    )

    return response.text.strip().replace("\n", "")
# 1. Generate title from document content
doc_title = generate_title_from_content(st.session_state.summary_output)

# 2. Create Firestore document
summary_doc_ref = summary_col.document()

summary_doc_ref.set({
    "title": doc_title,
    "created_at": firestore.SERVER_TIMESTAMP,
    "source": "summarizer"
})
summary_output1=st.session_state.summary_output
# 3. Store chunks
store_summary_chunks(summary_doc_ref,summary_output1 )

# 4. Track active summary
st.session_state.active_summary_id = summary_doc_ref.id
